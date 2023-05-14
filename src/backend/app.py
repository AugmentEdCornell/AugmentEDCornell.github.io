from pptx import Presentation
from flask import Flask, request
import os
import openai
from flask_cors import CORS

app = Flask(__name__)
CORS(app)

@app.route('/', methods=['POST'])
def chatGPT_Response():
    ppt_file = request.args.get('ppt_file')
    
    ppt_prompt = ""
    pres = Presentation(ppt_file)
    title = pres.core_properties.title

    ppt_prompt += "Presentation Name: " + title

    for num, slide in enumerate(pres.slides):
        if slide.shapes.title:
            heading = slide.shapes.title.text
            ppt_prompt += "\n --Slide Heading " + str(num)+ "--  :" + heading
            
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    ppt_prompt += paragraph.text + "\n"
                    
    print(ppt_prompt)

    openai.api_key = os.environ["OPENAI_API_KEY"]
    openai.organization = "org-otJnqs6IrdU7C4RhmoFFm0tY"

    engine = "text-davinci-003"
    text_prompt = "For students in grade 5. Generate a 3 question MCQ quiz with 4 options each. Also mention the correct answers for it at the end of all the questions. The quiz is on plant cell biology. \n\n"
    final_prompt = text_prompt  + ppt_prompt

    params = {
        "engine": engine,
        "prompt": final_prompt,
        "max_tokens": 300,
        "temperature": 0.5
    }

    response = openai.Completion.create(**params)

    if response.choices:
        result = response.choices[0].text
        print(result)
    else:
        print(f'Error: {response["error"]["message"]}')

    return result

if __name__ == "__main__":
    app.run()